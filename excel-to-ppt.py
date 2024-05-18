# The following program is code that I wrote to solve a problem that a press company faced. The problem is as follows:
# A very large powerpoint presentation must be given in order to award a large swath of journalists from many cities across the state of New Mexico.
# The data concerning the list of awards was compiled in an excel sheet, containing about 1330 cells of data that must be translated over to PowerPoint in a presentable fashion.
# It takes well over 100 hours to complete this process by hand, and many mistakes can crop up throughout the process.
# This program aims to cut the amount of hours and manpower required to create a slide from this massive amount of data down to mere minutes.

# The program simply grabs the data from a specified range of cells in a targeted Excel file, and spits that data out onto PowerPoint. Meanwhile, the program also edits the slides of the powerpoint to be more presentable.
# At present, the program takes about a minute or two to run and save a completed PowerPoint file.
# Some manual edits to image sizes on the powerpoint file may be required as it would be very difficult to scale every image perfectly on each slide, 
# not impossible per se, but about as time consuming as simply making the manual adjustments.
# The images are added to the slides using default size values.

# Refer to the pptx documentation for additional information on font editing and slide layouts with regards to this program.

# All libraries required for this code to run:
from tkinter import font
from openpyxl.workbook import Workbook
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

# The below code is required to initialize both the excel worksheet and the powerpoint files, as well as to ensure that images of any size can be manipulated with this program.
Image.MAX_IMAGE_PIXELS = None

workbook = load_workbook('NMPA Contest Winners 2023 For Power Point.xlsx')
worksheet = workbook.active

#Change the values in the brackets to specify the range of cells that need to be translated onto PowerPoint.
range = worksheet['A4':'E269']

presentation = Presentation()

# Default Function that writes and edits text for each slide, excluding the Intro slide
def set_text_and_color(paragraph, text, font_size, font_color):
    p = paragraph
    p.text = text
    font = p.font
    font.size = font_size
    font.color.rgb = font_color

def MakeIntroSlide():
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "New Mexico Press Association Awards Presentation"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)

    subtitle.text = "2023"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 0)


# This is the main function responsible for the creation of every slide after the intro slide, in this particular project, it creates over 250+ slides from over 1000 cells of data
def MakeSlides(range):
    i = 1 # Loop variable
    img_num = 2
    for cell in range:
        for x in cell:
            if i == 1:
                slide_register = presentation.slide_layouts[5]
                slide = presentation.slides.add_slide(slide_register)
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0, 0, 0)
                title = slide.shapes.title

                img_path = 'Contest Winners Art for Powerpont/Imgs/Renames/' + str(img_num) + '.jpg'

                # This exception must be placed in the code in order for images to be attached to each slide. 
                # I'm not exactly the sure the reasoning for the value error as the pptx library is built to work with images.
                # FileNotFoundError must also be placed as an exception in case certain slides simply do not have images attached to them.
                try:
                    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1), Inches(4), Inches(6))
                except (FileNotFoundError, ValueError):
                    pass
                img_num += 1

                title.text = x.value
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
                title.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

                i += 1
            elif i == 2:
                left = Inches(6.33)
                top = Inches(2)
                width = height = Inches(3.5)

                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True

                set_text_and_color(tf.add_paragraph(), x.value, Pt(24), RGBColor(255, 255, 0))

                tf.add_paragraph()

                i += 1
            elif i == 3:

                set_text_and_color(tf.add_paragraph(), x.value, Pt(28), RGBColor(255, 255, 0))

                tf.add_paragraph()
                
                i += 1
            elif i == 4:

                set_text_and_color(tf.add_paragraph(), x.value, Pt(28), RGBColor(255, 255, 0))

                tf.add_paragraph()

                i += 1
            elif i >= 5:

                set_text_and_color(tf.add_paragraph(), x.value, Pt(28), RGBColor(255, 255, 0))

                tf.add_paragraph()

                i = 1

MakeIntroSlide()

MakeSlides(range)

# We call presentation.save in order to actually create the powerpoint file! We can give it any name we would like in the ''
presentation.save('test.pptx')