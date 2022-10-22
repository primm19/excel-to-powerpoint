# This program is meant for converting cell data from Excel Files into readable text
# onto a PowerPoint project. This program also iterates through its directory to find
# and place .png files onto each slide. This program also has minor slide formatting
# espoused onto each slide. It is immensely useful in taking large spreadsheets of text
# and printing them out onto PowerPoint slides. WARNING: THIS PROGRAM WILL NOT RUN IF
# YOU HAVE POWERPOINT OPEN. YOU MUST HAVE IT CLOSED.

# Libraries -- I looked up Python-PPTX documentation a lot for this project
from tkinter import font
from openpyxl.workbook import Workbook
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path
from pptx.enum.text import PP_ALIGN

# Set the file name of the excel sheet here
wb = load_workbook('NMPA Contest Winners 2022 FINAL For Lee.xlsx')
ws = wb.active

# Specify your cell ranges here -- This may need to change depending on what cell ranges
# your cell data occupies across the spreadsheet. You may need to set separate ranges
# (as seen below) if you want to do entirely different things with different ranges.
range1 = ws['A2':'E307']
range2 = ws['B309': 'C310']

cellArrayLst = []

prs = Presentation()

# First Range Loop
i = 1 # We must iterate through a variable and a specified range for each number of cells that we want per slide
for cell in range1:# We must loop twice to get actual cell data, otherwise we are left with tuples which will not output to PowerPoint
    for x in cell:
        if i == 1:
            slide_register = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_register)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)
            title = slide.shapes.title
            title.text = x.value
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            # cellArrayLst.append(f'{x.value.upper()}')
            i += 1
        elif i == 2:
            left = Inches(6.33)
            top = Inches(2)
            width = height = Inches(3.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.text = x.value
            p.font.size = Pt(24)
            font = p.font
            font.color.rgb = RGBColor(255, 255, 0)
            p = tf.add_paragraph()
            # cellArrayLst.append(f'{x.value.upper()}')
            i += 1
        elif i == 3:
            p = tf.add_paragraph()
            p.text = x.value
            p.font.size = Pt(28)
            font = p.font
            font.color.rgb = RGBColor(255, 255, 0)
            p = tf.add_paragraph()
            # cellArrayLst.append(f'{x.value.upper()}')
            i += 1
        elif i == 4:
            p = tf.add_paragraph()
            p.text = x.value
            font = p.font
            font.color.rgb = RGBColor(255, 255, 0)
            p = tf.add_paragraph()
            # cellArrayLst.append(f'{x.value.upper()}')
            i += 1
        elif i >= 5:
            p = tf.add_paragraph()
            p.text = x.value
            font = p.font
            font.color.rgb = RGBColor(255, 255, 0)
            p = tf.add_paragraph()
            # cellArrayLst.append(f'{x.value.upper()}')
            # print(cellArrayLst)
            # Need to pass array through function here
            # cellArrayLst = []
            i = 1
        # An else statement like the one below is written simply for the developer to understand
        # that they may have created an infinite loop or not properly created the number of
        # iterations they wanted to.
        # WARNING -- It is still possible to make a mistake here and have the wrong number of
        # iterations, especially if you want a fewer number of cells per slide than
        # the current established range of iterations. You simply must change the number of 'if'
        # or 'elif' statements.
        else:
            print("Error: Counter('i') has exceeded max limit, please set a new limit or change the your worksheet range in Python.")

# Range Two Loop
j = 1
for cell in range2:
    for x in cell:
        if j == 1:
            slide_register = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_register)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)
            title = slide.shapes.title
            title.text = x.value
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
            j += 1
        elif j >= 2:
            left = Inches(2.25)
            top = Inches(2)
            width = height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = x.value
            p.font.size = Pt(24)
            font = p.font
            font.color.rgb = RGBColor(255, 255, 0)
            p = tf.add_paragraph()
            j = 1
        else:
            print("Error: Counter('j') has exceeded max limit, please set a new limit or change the your worksheet range in Python.")

# Final PowerPoint file -- Change this if you want to change the powerpoint filename.
prs.save('New Mexico Press Association Awards Presentation.pptx')
