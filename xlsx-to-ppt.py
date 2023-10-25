from tkinter import font
from openpyxl.workbook import Workbook
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

workbook = load_workbook('2023 NNAF Contest Winners ALL fields kd edit.xlsx')
worksheet = workbook.active

advertiser_range = worksheet['B2':'D4'] #complete
almanance_range = worksheet['B5':'D13'] #complete
antelope_range = worksheet['B14':'D30'] #complete
ark_range = worksheet['B31':'D53'] #complete
brentwood_range = worksheet['B54':'D58'] #complete
buffalo_range = worksheet['B59':'D75'] #complete
casa_range = worksheet['B76':'D90'] #complete
cedar_range = worksheet['B91':'D100'] 
chronicle_range = worksheet['B101':'D102'] 
citizen_range = worksheet['B103':'D104']
cody_range = worksheet['B105':'D108']
colbert_range = worksheet['B109':'D109']
custer_range = worksheet['B110':'D110']
washingtonlaw_range = worksheet['B111':'D111']
douglas_range = worksheet['B112':'D116']
edgerton_range = worksheet['B117':'D127']
eldon_range = worksheet['B128':'D129']
eldridge_range = worksheet['B130':'D137']
examiner_range = worksheet['B138':'D150']
farm_range = worksheet['B151':'D152']
ftbend_range = worksheet['B153':'D155']
gazebo_range = worksheet['B156':'D160']
greenvilledaily_range = worksheet['B161':'D161']
greenville_range = worksheet['B162':'D162']
henderson_range = worksheet['B163':'D167']
highlands_range = worksheet['B168':'D183']
hinsdalean_range = worksheet['B184':'D187']
hyattsville_range = worksheet['B188':'D190']
idahomtn_range = worksheet['B191':'D199']
jackson_range = worksheet['B200':'D215']
jefferson_range = worksheet['B216':'D216']
journal_range = worksheet['B217':'D225']
journalnews_range = worksheet['B226':'D232']
lahontan_range = worksheet['B233':'D235']
laurel_range = worksheet['B236':'D237']
leelanau_range = worksheet['B238':'D263']
madison_range = worksheet['B264':'D264']
manchester_range = worksheet['B265':'D267']
marietta_range = worksheet['B268':'D274']
midhudson_range = worksheet['B275':'D275']
moonshine_range = worksheet['B276':'D286']
myrtle_range = worksheet['B287':'D288']
nwiowa_range = worksheet['B289':'D296']
newsletter_range = worksheet['B297':'D311']
newsgazette_range = worksheet['B312':'D315']
northlight_range = worksheet['B316':'D318']
oakland_range = worksheet['B319':'D321']
observer_range = worksheet['B322':'D323']
observertrib_range = worksheet['B324':'D324']
omak_range = worksheet['B325':'D325']
ozona_range = worksheet['B326':'D327']
parkcities_range = worksheet['B328':'D334']
petersburg_range = worksheet['B335':'D336']
phillygay_range = worksheet['B337':'D345']
photonews_range = worksheet['B346':'D346']
pike_range = worksheet['B347':'D347']
pilot_range = worksheet['B348':'D355']
portarkansas_range = worksheet['B356':'D359']
preston_range = worksheet['B360':'D365']
redhook_range = worksheet['B366':'D367']
rockcounty_range = worksheet['B368':'D368']
sanfernando_range = worksheet['B369':'D370']
shelter_range = worksheet['B371':'D373']
sioux_range = worksheet['B374':'D393']
southwest_range = worksheet['B394':'D408']
swcollege_range = worksheet['B409':'D420']
sparta_range = worksheet['B421':'D427']
spartamonroe_range = worksheet['B428':'D428']
standardbanner_range = worksheet['B429':'D435']
stanton_range = worksheet['B436':'D444']
stateport_range = worksheet['B445':'D447']
steele_range = worksheet['B448':'D460']
tampa_range = worksheet['B461':'D463']
taos_range = worksheet['B464':'D511']
maine_range = worksheet['B512':'D516']
times_range = worksheet['B517':'D517']
tioga_range = worksheet['B518':'D521']
township_range = worksheet['B522':'D526']
trempealeau_range = worksheet['B527':'D527']
tuscola_range = worksheet['B528':'D529']
uinta_range = worksheet['B530':'D532']
uvalde_range = worksheet['B533':'D536']
vilas_range = worksheet['B537':'D543']
village_range = worksheet['B544':'D548']
wallkill_range = worksheet['B549':'D549']
warwick_range = worksheet['B550':'D553']
watertown_range = worksheet['B554':'D555']
wilton_range = worksheet['B556':'D556']
wright_range = worksheet['B557':'D563']
wyoming_range = worksheet['B564':'D581']
wyomingtruth_range = worksheet['B582':'D586']
yankton_range = worksheet['B587':'D605']

presentation = Presentation()

def MakeSlidesNoImages(range, header_name, presentation, text_top, header_top, text_left, header_left, text_width, header_height, text_height, header_width):
    i = 1
    for cell in range:
        for x in cell:
            if i == 1:
                header_path = 'Headers/' + header_name + '.jpg'
                blank_slide = presentation.slide_layouts[6]
                slide = presentation.slides.add_slide(blank_slide)
                slide.shapes.add_picture(header_path, header_left, header_top, header_width, header_height)
            # try:
            #     slide.shapes.add_picture(img_path, img_left, img_top, img_height, img_width)
            # except (FileNotFoundError, ValueError):
            #     pass
                txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.add_paragraph()
                p.text = x.value
                p.alignment = PP_ALIGN.CENTER
                i += 1
            # img_num += 1
            elif i == 2:
                p = tf.add_paragraph()
                p.text = x.value
                p.alignment = PP_ALIGN.CENTER
                i += 1
            elif i >= 3:
                p = tf.add_paragraph()
                p.text = x.value
                p.alignment = PP_ALIGN.CENTER
                i = 1

def MakeSlidesWithImages(range, img_num, header_name, presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height):
    i = 1
    img_num = img_num
    for cell in range:
        for x in cell:
            if i == 1:
                header_path = 'Headers/' + header_name + '.jpg'
                img_path = 'Images/' + str(img_num) + '.jpg'
                blank_slide = presentation.slide_layouts[6]
                slide = presentation.slides.add_slide(blank_slide)
                slide.shapes.add_picture(header_path, header_left, header_top, header_width, header_height)
                try:
                    slide.shapes.add_picture(img_path, img_left, img_top, img_height, img_width)
                except (FileNotFoundError, ValueError):
                    pass
                txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.add_paragraph()
                p.text = x.value
                i += 1
                img_num += 1
            elif i == 2:
                p = tf.add_paragraph()
                p.text = x.value
                i += 1
            elif i >= 3:
                p = tf.add_paragraph()
                p.text = x.value
                i = 1

#Advertiser
header_name = ""
text_top = Inches(3)
header_top = Inches(0.5)
text_left = Inches(3)
header_left = Inches(1.75)
text_width = Inches(5)
header_height = Inches(2)
text_height = Inches(3)
header_width = Inches(6.4)

MakeSlidesNoImages(advertiser_range, "Advertiser", presentation, text_top, header_top, text_left, header_left, text_width, header_height, text_height, header_width)

#Almanance 
text_left = Inches(1)
text_top = Inches(2.5)
text_width = Inches(4)
text_height = Inches(3)
img_top = Inches(1)
img_left = Inches(5.5)
img_height = Inches(4)
img_width = Inches(6)
header_width = Inches(5)
header_left = Inches(0.5)
header_top = Inches(0.5)
header_height = Inches(1)

MakeSlidesWithImages(almanance_range, 4, "Almanance", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Antelope
header_width = Inches(4.5)
header_height = Inches(1.5)
text_top = Inches(3)

MakeSlidesWithImages(antelope_range, 13, "Antelope", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Ark
header_height = Inches(2.09)
header_width = Inches(4.19)
header_top = Inches(0)
img_top = Inches(0.5)

MakeSlidesWithImages(ark_range, 30, "Ark", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Brentwood
header_height = Inches(1)
header_width = Inches(3.5)
header_top = Inches(0.75)
header_left = Inches(1)

MakeSlidesWithImages(brentwood_range, 53, "Brentwood", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Buffalo
header_height = Inches(1)
header_width = Inches(6)
header_left = Inches(0.5)
header_top = Inches(0.25)
img_top = Inches(1.33)

MakeSlidesWithImages(buffalo_range, 58, "Buffalo", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Casa
header_height = Inches(1.5)
header_width = Inches(4.5)
img_top = Inches(0.75)

MakeSlidesWithImages(casa_range, 75, "Casa", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Cedar
header_height = Inches(1)
header_width = Inches(6)
img_top = Inches(1.33)

MakeSlidesWithImages(cedar_range, 90, "Cedar", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Chronicle
header_height = Inches(1.5)
header_width = Inches(4)
header_top = Inches(1)

MakeSlidesWithImages(chronicle_range, 100, "Chronicle", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Citizen
header_height = Inches(1)
header_width = Inches(5)
header_top = Inches(0.5)

MakeSlidesWithImages(citizen_range, 102, "Citizen", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Cody
header_height = Inches(1)
header_width = Inches(6)
header_top = Inches(0.5)

MakeSlidesWithImages(cody_range, 104, "Cody", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Colbert
header_height = Inches(1)
header_width = Inches(6)
header_top = Inches(1.5)
header_left = Inches(2)
text_left = Inches(3.5)

MakeSlidesWithImages(colbert_range, 108, "Colbert", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Custer
header_height = Inches(4)
header_width = Inches(6)
header_top = Inches(0.5)
text_top = Inches(4.5)

MakeSlidesWithImages(custer_range, 109, "Custer", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Washington Law
header_height = Inches(1.5)
header_width = Inches(2.5)
header_top = Inches(0.5)
header_left = Inches(1)
text_left = Inches(1)
text_top = Inches(2.5)
img_top = Inches(1)
img_left = Inches(5)

MakeSlidesWithImages(washingtonlaw_range, 110, "WashingtonLaw", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Douglas
header_height = Inches(1)
header_width = Inches(6)
img_left = Inches(5.5)
header_top = Inches(0)
header_left = Inches(0.25)

MakeSlidesWithImages(douglas_range, 111, "Douglas", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Edgerton
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(edgerton_range, 116, "Edgerton", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Eldon
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(eldon_range, 127, "Eldon", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Eldridge
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(eldridge_range, 129, "Eldridge", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Examiner
header_height = Inches(1.5)
header_width = Inches(6)
header_left = Inches(0)
img_left = Inches(5.5)
img_top = Inches(1.5)
text_top = Inches(3)

MakeSlidesWithImages(examiner_range, 137, "Examiner", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Farm
header_height = Inches(1.3)
header_width = Inches(4.75)
img_top = Inches(1)
img_left = Inches(5.5)
header_top = Inches(0.5)
header_left = Inches(0.5)

MakeSlidesWithImages(farm_range, 150, "Farm", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Ft. Bend
header_height = Inches(1)
header_width = Inches(6)
header_top = Inches(0)
header_left = Inches(0)

MakeSlidesWithImages(ftbend_range, 152, "FtBend", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Gazebo
header_height = Inches(1.3)
header_width = Inches(4.75)

MakeSlidesWithImages(gazebo_range, 155, "Gazebo", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Greenville Daily
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(greenvilledaily_range, 160, "GreenvilleDaily", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Greenville
header_height = Inches(1)
header_width = Inches(6)
img_top = Inches(1.5)

MakeSlidesWithImages(greenville_range, 161, "Greenville", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Henderson
header_height = Inches(1)
header_width = Inches(6)
img_top = Inches(1)

MakeSlidesWithImages(henderson_range, 162, "Henderson", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Highlands
header_height = Inches(1.8)
header_width = Inches(4.8)

MakeSlidesWithImages(highlands_range, 167, "Highlands", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Hinsdalean
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(hinsdalean_range, 183, "Hinsdalean", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Hyattsville
header_height = Inches(1)
header_width = Inches(5)

MakeSlidesWithImages(hyattsville_range, 187, "Hyattsville", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Idaho Mountain
header_height = Inches(2.6)
header_width = Inches(4.7)

MakeSlidesWithImages(idahomtn_range, 190, "IdahoMtn", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Jackson
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(jackson_range, 199, "Jackson", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Jefferson
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(jefferson_range, 215, "Jefferson", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Journal
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(journal_range, 216, "Journal", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Journal News
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(journalnews_range, 225, "JournalNews", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Lahontan
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(lahontan_range, 232, "Lahontan", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Laurel
header_height = Inches(1.5)
header_width = Inches(4)

MakeSlidesWithImages(laurel_range, 235, "Laurel", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Leelanau
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(leelanau_range, 237, "Leelanau", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Madison
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(madison_range, 263, "Madison", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Manchester
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(manchester_range, 264, "Manchester", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Marietta
header_height = Inches(1.5)
header_width = Inches(4)

MakeSlidesWithImages(marietta_range, 267, "Marietta", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Midhudson
header_height = Inches(1.5)
header_width = Inches(4)

MakeSlidesWithImages(midhudson_range, 274, "Midhudson", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Moonshine Ink
header_height = Inches(2.5)
header_width = Inches(3)

MakeSlidesWithImages(moonshine_range, 275, "Moonshine", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Myrtle
header_height = Inches(2)
header_width = Inches(4)

MakeSlidesWithImages(myrtle_range, 286, "Myrtle", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#NW Iowa
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(nwiowa_range, 288, "NWIowa", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Newsletter
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(newsletter_range, 296, "Newsletter", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Newsgazette
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(newsgazette_range, 311, "Newsgazette", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Northlight
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(northlight_range, 315, "Northlight", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Oakland
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(oakland_range, 318, "Oakland", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Observer
header_height = Inches(1)
header_width = Inches(10)

MakeSlidesWithImages(observer_range, 321, "Observer", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Observer Tribune
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(observertrib_range, 323, "ObserverTribune", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Omak
header_height = Inches(1)
header_width = Inches(4.5)

MakeSlidesWithImages(omak_range, 324, "Omak", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Ozona
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(ozona_range, 325, "Ozona", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Park Cities
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(parkcities_range, 327, "ParkCities", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Petersburg
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(petersburg_range, 334, "Petersburg", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Philidelphia Gay News
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(phillygay_range, 336, "PhillyGay", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Photo News
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(photonews_range, 345, "PhotoNews", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Pike
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(pike_range, 346, "Pike", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Pilot
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(pilot_range, 347, "Pilot", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Port Aransas
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(portarkansas_range, 355, "PortArkansas", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Preston
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(preston_range, 359, "Preston", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Red Hook
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(redhook_range, 365, "RedHook", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Rock County
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(rockcounty_range, 367, "RockCounty", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#San Fernando
header_height = Inches(2)
header_width = Inches(4)

MakeSlidesWithImages(sanfernando_range, 368, "SanFernando", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Shelter
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(shelter_range, 370, "Shelter", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Sioux
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(sioux_range, 373, "Sioux", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Southwest
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(southwest_range, 393, "Southwest", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Southwest College
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(swcollege_range, 408, "SWCollege", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Sparta
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(sparta_range, 420, "Sparta", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Sparta-Monroe
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(spartamonroe_range, 427, "SpartaMonroe", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Standard Banner
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(standardbanner_range, 428, "StandardBanner", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Stanton
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(stanton_range, 435, "Stanton", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Stateport
header_height = Inches(2)
header_width = Inches(3.5)

MakeSlidesWithImages(stateport_range, 444, "Stateport", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Steele
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(steele_range, 447, "Steele", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Tampa Bay
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(tampa_range, 460, "Tampa", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Taos
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(taos_range, 463, "Taos", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Maine
header_height = Inches(1.5)
header_width = Inches(3)

MakeSlidesWithImages(maine_range, 511, "Maine", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Times
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(times_range, 516, "Times", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Tioga
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(tioga_range, 517, "Tioga", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Township
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(township_range, 521, "Township", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Trempealeau
header_height = Inches(1.5)
header_width = Inches(3.5)

MakeSlidesWithImages(trempealeau_range, 526, "Trempealeau", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Tuscola
header_height = Inches(1.9)
header_width = Inches(4.5)

MakeSlidesWithImages(tuscola_range, 527, "Tuscola", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Uinta
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(uinta_range, 529, "Uinta", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Uvalde
header_height = Inches(0.75)
header_width = Inches(6)

MakeSlidesWithImages(uvalde_range, 532, "Uvalde", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Vilas
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(vilas_range, 536, "Vilas", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Village
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(village_range, 543, "Village", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Wallkill
header_height = Inches(1.5)
header_width = Inches(4.5)

MakeSlidesWithImages(wallkill_range, 548, "Wallkill", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Warwick
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(warwick_range, 549, "Warwick", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Watertown
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(watertown_range, 553, "Watertown", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Wilton
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(wilton_range, 555, "Wilton", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Wright
header_height = Inches(1.5)
header_width = Inches(4.25)

MakeSlidesWithImages(wright_range, 556, "Wright", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Wyoming
header_height = Inches(1.75)
header_width = Inches(4.5)

MakeSlidesWithImages(wyoming_range, 563, "Wyoming", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Wyoming Truth
header_height = Inches(2.75)
header_width = Inches(2.75)

MakeSlidesWithImages(wyomingtruth_range, 581, "WyomingTruth", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

#Yankton
header_height = Inches(1)
header_width = Inches(6)

MakeSlidesWithImages(yankton_range, 586, "Yankton", presentation, img_top, img_left, img_height, img_width, header_width, header_left, header_top, header_height, text_left, text_top, text_width, text_height)

presentation.save('test.pptx')
